# -*- coding: utf-8 -*-
"""
模板引擎 - PPT 模板加载与占位符填充
支持动态占位符配置 + 自动扫描生成
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import json
import re
from typing import Dict, List, Optional, Tuple
import logging
from datetime import datetime

logger = logging.getLogger(__name__)


class TemplateEngine:
    """PPT 模板引擎 - 支持占位符替换 + 自动扫描"""
    
    def __init__(self, base_dir: str = None, auto_scan: bool = True):
        """
        初始化模板引擎
        
        Args:
            base_dir: 基础目录
            auto_scan: 如果 placeholders.json 不存在，是否自动从模板扫描生成
        """
        # 支持 EXE 打包：优先使用 EXE_WORK_DIR 环境变量
        self.base_dir = os.environ.get('EXE_WORK_DIR') or base_dir or os.path.dirname(os.path.dirname(__file__))
        self.templates_dir = os.path.join(self.base_dir, 'templates')
        self.artifacts_dir = os.path.join(self.base_dir, 'artifacts')
        self.placeholders_file = os.path.join(self.artifacts_dir, 'placeholders.json')
        
        # 确保目录存在
        os.makedirs(self.templates_dir, exist_ok=True)
        
        # 加载或生成占位符配置
        self.placeholders = self._load_or_scan_placeholders(auto_scan)
    
    def _load_or_scan_placeholders(self, auto_scan: bool) -> Dict:
        """加载占位符配置，如果不存在则自动扫描"""
        if os.path.exists(self.placeholders_file):
            logger.info(f"加载占位符配置：{self.placeholders_file}")
            with open(self.placeholders_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        
        if auto_scan:
            logger.info("占位符配置不存在，尝试从模板扫描生成...")
            default_template = os.path.join(self.templates_dir, '销售分析报告_标准模板.pptx')
            if os.path.exists(default_template):
                placeholders = self.scan_placeholders_from_template(default_template)
                self._save_placeholders(placeholders)
                logger.info(f"已自动生成占位符配置：{self.placeholders_file}")
                return placeholders
        
        logger.warning("未找到占位符配置，使用空配置")
        return self._get_default_placeholders()
    
    def _get_default_placeholders(self) -> Dict:
        """返回默认占位符配置（v3.0 完全配置化版）"""
        return {
            "version": "3.0",
            "description": "销售分析报告 PPT 占位符定义（完全配置化版）",
            "generated_at": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "template_file": "销售分析报告_标准模板.pptx",
            "placeholders": {
                "charts": {
                    "CHART:sales_by_person": {
                        "description": "销售员业绩横向条形图",
                        "data_source": "销售员业绩",
                        "chart_type": "bar_horizontal",
                        "x_field": "总销售额",
                        "y_field": "销售员",
                        "title": "销售员业绩表现分析",
                        "slide_index": 3
                    },
                    "CHART:product_pie": {
                        "description": "产品销售占比环形图",
                        "data_source": "产品占比",
                        "chart_type": "pie",
                        "category_field": "产品",
                        "value_field": "占比",
                        "title": "产品销售结构与占比",
                        "figsize": [6, 6],
                        "slide_index": 4
                    },
                    "CHART:city_ranking": {
                        "description": "城市销售业绩排名",
                        "data_source": "城市排名",
                        "chart_type": "bar_horizontal",
                        "x_field": "总销售额",
                        "y_field": "城市",
                        "title": "福建各城市销售业绩排名",
                        "slide_index": 5
                    },
                    "CHART:customer_comparison": {
                        "description": "新老客销售对比",
                        "data_source": "客户类型",
                        "chart_type": "column_clustered",
                        "category_field": "客户属性",
                        "series": ["总销售额", "订单数", "客单价"],
                        "title": "新老客销售贡献对比",
                        "slide_index": 6
                    },
                    "CHART:monthly_trend": {
                        "description": "月度销售趋势",
                        "data_source": "月度趋势",
                        "chart_type": "line",
                        "x_field": "年月",
                        "y_field": "总销售额",
                        "title": "销售趋势分析",
                        "slide_index": 7
                    },
                    "CHART:heatmap": {
                        "description": "销售员 - 产品热力图",
                        "data_source": "销售员 - 产品",
                        "chart_type": "heatmap",
                        "index_field": "销售员",
                        "columns": null,
                        "title": "销售员 - 产品销售能力矩阵分析",
                        "slide_index": 8
                    }
                },
                "insights": {
                    "INSIGHT:kpi_summary": {"description": "核心指标洞察", "slide_index": 2, "skill": "data-insight", "page": 4},
                    "INSIGHT:sales_by_person": {"description": "销售员业绩洞察", "slide_index": 3, "skill": "data-insight", "page": 5},
                    "INSIGHT:product": {"description": "产品结构洞察", "slide_index": 4, "skill": "data-insight", "page": 6},
                    "INSIGHT:city": {"description": "城市区域洞察", "slide_index": 5, "skill": "data-insight", "page": 7},
                    "INSIGHT:customer": {"description": "客户类型洞察", "slide_index": 6, "skill": "data-insight", "page": 8},
                    "INSIGHT:trend": {"description": "时间趋势洞察", "slide_index": 7, "skill": "data-insight", "page": 9},
                    "INSIGHT:heatmap": {"description": "销售员 - 产品矩阵洞察", "slide_index": 8, "skill": "data-insight", "page": 10},
                    "INSIGHT:abnormal": {"description": "异常订单洞察", "slide_index": 9, "skill": "data-insight", "page": 11},
                    "INSIGHT:conclusion": {"description": "核心结论", "slide_index": 10, "skill": "data-insight", "page": 12},
                    "INSIGHT:strategy": {"description": "落地策略", "slide_index": 11, "skill": "data-insight", "page": 13}
                },
                "tables": {
                    "TABLE:abnormal_orders": {
                        "description": "异常订单明细表",
                        "data_source": "异常订单",
                        "columns": null,
                        "slide_index": 9
                    }
                },
                "text": {
                    "TEXT:report_title": {"description": "报告标题", "default": "销售数据分析报告", "slide_index": 0},
                    "TEXT:report_subtitle": {"description": "报告副标题", "default": "业绩复盘 · 洞察归因 · 策略建议", "slide_index": 0},
                    "TEXT:report_date": {"description": "报告日期", "format": "YYYY-MM-DD", "slide_index": 0},
                    "KPI:cards": {"description": "核心 KPI 卡片内容", "default": "", "slide_index": 2}
                }
            },
            "slide_mapping": {
                "0": {"name": "封面页", "layout": "title"},
                "1": {"name": "目录页", "layout": "toc"},
                "2": {"name": "核心指标", "layout": "kpi"},
                "3": {"name": "销售员业绩", "layout": "chart_text"},
                "4": {"name": "产品结构", "layout": "chart_text"},
                "5": {"name": "城市排名", "layout": "chart_text"},
                "6": {"name": "客户类型", "layout": "chart_text"},
                "7": {"name": "时间趋势", "layout": "chart_text"},
                "8": {"name": "销售员 - 产品", "layout": "chart_text"},
                "9": {"name": "异常订单", "layout": "table_text"},
                "10": {"name": "核心结论", "layout": "content"},
                "11": {"name": "落地策略", "layout": "content"},
                "12": {"name": "结束页", "layout": "title"}
            },
            "report_settings": {
                "date_range_auto_detect": true,
                "date_range_format": "统计周期：{start} - {end} | 日期：{today}"
            }
        }
    
    def _save_placeholders(self, placeholders: Dict):
        """保存占位符配置"""
        with open(self.placeholders_file, 'w', encoding='utf-8') as f:
            json.dump(placeholders, f, ensure_ascii=False, indent=2)
    
    def scan_placeholders_from_template(self, template_path: str) -> Dict:
        """
        从 PPT 模板扫描占位符
        
        Args:
            template_path: PPT 模板文件路径
        
        Returns:
            占位符配置字典
        """
        logger.info(f"扫描模板占位符：{template_path}")
        prs = Presentation(template_path)
        
        placeholders = {
            "version": "1.0",
            "description": "从模板自动扫描生成",
            "generated_at": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "template_file": os.path.basename(template_path),
            "placeholders": {
                "charts": {},
                "insights": {},
                "tables": {},
                "text": {}
            },
            "slide_mapping": {}
        }
        
        # 正则表达式匹配占位符（兼容单方括号和双方括号）
        chart_pattern = re.compile(r'\[\[?CHART:(\w+)\]\]?')
        insight_pattern = re.compile(r'\{\{INSIGHT:(\w+)\}\}')
        table_pattern = re.compile(r'\[\[?TABLE:(\w+)\]\]?')
        text_pattern = re.compile(r'\[\[?TEXT:(\w+)\]\]?')
        kpi_pattern = re.compile(r'\[\[?KPI:(\w+)\]\]?')
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_info = {
                "name": f"Slide {slide_idx}",
                "layout": "content"
            }
            
            for shape_idx, shape in enumerate(slide.shapes):
                # 检查 shape 名称（图表占位符）
                if hasattr(shape, 'name') and shape.name and shape.name.startswith('CHART:'):
                    chart_key = shape.name.split(':')[1]
                    placeholders["placeholders"]["charts"][f"CHART:{chart_key}"] = {
                        "description": f"图表：{chart_key}",
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "by_name": True
                    }
                    logger.info(f"  发现图表占位符 (by name): CHART:{chart_key} (slide {slide_idx})")
                
                # 检查文本框内容
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        
                        # 匹配各种占位符
                        for pattern, ptype, pkey_prefix in [
                            (chart_pattern, 'charts', 'CHART:'),
                            (insight_pattern, 'insights', 'INSIGHT:'),
                            (table_pattern, 'tables', 'TABLE:'),
                            (text_pattern, 'text', 'TEXT:'),
                            (kpi_pattern, 'text', 'KPI:'),
                        ]:
                            match = pattern.search(text)
                            if match:
                                key = match.group(1)
                                full_key = f"{pkey_prefix}{key}"
                                
                                if full_key not in placeholders["placeholders"][ptype]:
                                    placeholders["placeholders"][ptype][full_key] = {
                                        "description": f"{ptype}：{key}",
                                        "slide_index": slide_idx,
                                        "default": ""
                                    }
                                    logger.info(f"  发现{ptype}占位符：{full_key} (slide {slide_idx})")
            
            placeholders["slide_mapping"][str(slide_idx)] = slide_info
        
        # 推断数据源（基于占位符名称）
        self._infer_data_sources(placeholders)
        
        return placeholders
    
    def _infer_data_sources(self, placeholders: Dict):
        """根据占位符名称推断数据源"""
        # 图表数据源映射
        chart_data_map = {
            'sales_by_person': '销售员业绩',
            'product_pie': '产品占比',
            'city_ranking': '城市排名',
            'customer_comparison': '客户类型',
            'monthly_trend': '月度趋势',
            'heatmap': '销售员 - 产品',
            'quarterly_comparison': '季度对比',
        }
        
        for chart_key in placeholders["placeholders"]["charts"].keys():
            key_name = chart_key.split(':')[1]
            if key_name in chart_data_map:
                placeholders["placeholders"]["charts"][chart_key]["data_source"] = chart_data_map[key_name]
        
        # 洞察页码映射
        insight_pages = [
            'kpi_summary', 'sales_by_person', 'product', 'city', 'customer',
            'trend', 'heatmap', 'abnormal', 'conclusion', 'strategy'
        ]
        for i, insight_key in enumerate(insight_pages):
            full_key = f"INSIGHT:{insight_key}"
            if full_key in placeholders["placeholders"]["insights"]:
                placeholders["placeholders"]["insights"][full_key]["page"] = i + 4  # 从第 4 页开始
                placeholders["placeholders"]["insights"][full_key]["skill"] = "data-insight"
    
    def _load_placeholders(self) -> Dict:
        """加载占位符配置（不自动扫描）"""
        if os.path.exists(self.placeholders_file):
            with open(self.placeholders_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        return self._get_default_placeholders()
    
    def load_template(self, template_name: str) -> Presentation:
        """加载 PPT 模板"""
        template_path = os.path.join(self.templates_dir, template_name)
        if not os.path.exists(template_path):
            logger.info(f"模板不存在，创建基础模板：{template_path}")
            prs = self._create_default_template(template_path)
            # 新模板创建后，重新扫描占位符
            self.placeholders = self.scan_placeholders_from_template(template_path)
            return prs
        
        logger.info(f"加载模板：{template_path}")
        return Presentation(template_path)
    
    def _create_default_template(self, output_path: str) -> Presentation:
        """创建默认模板（带占位符）"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        colors = {
            'primary': RGBColor(31, 78, 121),
            'secondary': RGBColor(46, 117, 182),
            'accent': RGBColor(231, 76, 60),
            'text': RGBColor(44, 62, 80),
            'light': RGBColor(236, 240, 241),
        }
        
        # 第 1 页：封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[TEXT:report_title]"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = colors['primary']
        p.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1.5))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[TEXT:report_subtitle]"
        p.font.size = Pt(24)
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.CENTER
        
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[TEXT:report_date]"
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 第 2 页：目录
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "报告目录"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for item in ['1. 分析背景', '2. 核心指标', '3. 多维度分析', '4. 专项洞察', '5. 结论策略', '6. 结束页']:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
        
        # 第 3 页：核心指标
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "核心经营指标总览"
        kpi_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(4))
        kpi_box.name = "KPI:cards"
        tf = kpi_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "[KPI:cards]"
        p.font.size = Pt(24)
        p.font.color.rgb = colors['text']
        p.line_spacing = 1.5
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        
        insight_box = slide.shapes.add_textbox(Inches(5.8), Inches(1.5), Inches(7.033), Inches(5))
        tf = insight_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "{{INSIGHT:kpi_summary}}"
        p.font.size = Pt(14)
        p.line_spacing = 1.5
        p.space_after = Pt(12)
        
        # 第 4 页：销售员业绩
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "销售员业绩表现分析"
        chart_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(8), Inches(4.5)
        )
        chart_shape.fill.solid()
        chart_shape.fill.fore_color.rgb = colors['light']
        chart_shape.line.color.rgb = colors['secondary']
        chart_shape.text_frame.text = "[CHART:sales_by_person]"
        chart_shape.name = "CHART:sales_by_person"
        insight_box = slide.shapes.add_textbox(Inches(8.8), Inches(1.5), Inches(4), Inches(4.5))
        insight_box.text_frame.paragraphs[0].text = "{{INSIGHT:sales_by_person}}"
        
        # 第 5 页：产品结构
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "产品销售结构与占比分析"
        chart_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5)
        )
        chart_shape.fill.solid()
        chart_shape.fill.fore_color.rgb = colors['light']
        chart_shape.line.color.rgb = colors['secondary']
        chart_shape.text_frame.text = "[CHART:product_pie]"
        chart_shape.name = "CHART:product_pie"
        insight_box = slide.shapes.add_textbox(Inches(6.3), Inches(1.5), Inches(6.5), Inches(5.5))
        insight_box.text_frame.paragraphs[0].text = "{{INSIGHT:product}}"
        
        # 第 6 页：城市排名
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "福建各城市销售业绩排名"
        chart_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(8), Inches(4.5)
        )
        chart_shape.fill.solid()
        chart_shape.fill.fore_color.rgb = colors['light']
        chart_shape.line.color.rgb = colors['secondary']
        chart_shape.text_frame.text = "[CHART:city_ranking]"
        chart_shape.name = "CHART:city_ranking"
        insight_box = slide.shapes.add_textbox(Inches(8.8), Inches(1.5), Inches(4), Inches(4.5))
        insight_box.text_frame.paragraphs[0].text = "{{INSIGHT:city}}"
        
        # 第 7 页：客户类型
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "新老客销售贡献对比分析"
        chart_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(8), Inches(4.5)
        )
        chart_shape.fill.solid()
        chart_shape.fill.fore_color.rgb = colors['light']
        chart_shape.line.color.rgb = colors['secondary']
        chart_shape.text_frame.text = "[CHART:customer_comparison]"
        chart_shape.name = "CHART:customer_comparison"
        insight_box = slide.shapes.add_textbox(Inches(8.8), Inches(1.5), Inches(4), Inches(4.5))
        insight_box.text_frame.paragraphs[0].text = "{{INSIGHT:customer}}"
        
        # 第 8 页：月度趋势
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "2022 年上半年销售趋势分析"
        chart_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12), Inches(4)
        )
        chart_shape.fill.solid()
        chart_shape.fill.fore_color.rgb = colors['light']
        chart_shape.line.color.rgb = colors['secondary']
        chart_shape.text_frame.text = "[CHART:monthly_trend]"
        chart_shape.name = "CHART:monthly_trend"
        insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(12.333), Inches(1.3))
        insight_box.text_frame.paragraphs[0].text = "{{INSIGHT:trend}}"
        
        # 第 9 页：销售员 - 产品热力图
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "销售员 - 产品销售能力矩阵分析"
        chart_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12), Inches(4)
        )
        chart_shape.fill.solid()
        chart_shape.fill.fore_color.rgb = colors['light']
        chart_shape.line.color.rgb = colors['secondary']
        chart_shape.text_frame.text = "[CHART:heatmap]"
        chart_shape.name = "CHART:heatmap"
        insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(12.333), Inches(1.3))
        insight_box.text_frame.paragraphs[0].text = "{{INSIGHT:heatmap}}"
        
        # 第 10 页：异常订单表格
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "异常订单专项排查与分析"
        table_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(3.5))
        table_box.text_frame.paragraphs[0].text = "[TABLE:abnormal_orders]"
        table_box.name = "TABLE:abnormal_orders"
        insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.333), Inches(1.5))
        insight_box.text_frame.paragraphs[0].text = "{{INSIGHT:abnormal}}"
        
        # 第 11 页：核心结论
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "全量分析核心洞察结论"
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = "{{INSIGHT:conclusion}}"
        tf.paragraphs[0].line_spacing = 1.2
        tf.paragraphs[0].space_after = Pt(12)
        
        # 第 12 页：落地策略
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "可落地销售优化策略建议"
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = "{{INSIGHT:strategy}}"
        tf.paragraphs[0].line_spacing = 1.2
        tf.paragraphs[0].space_after = Pt(12)
        
        # 第 13 页：结束页
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "感谢观看"
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[TEXT:report_date]"
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.CENTER
        
        prs.save(output_path)
        logger.info(f"默认模板已创建：{output_path}")
        return prs
    
    def replace_text_placeholders(self, prs: Presentation, replacements: Dict) -> int:
        """替换文本占位符"""
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for placeholder, value in replacements.items():
                            if placeholder in paragraph.text:
                                value_str = '\n'.join(value) if isinstance(value, list) else str(value)
                                paragraph.text = paragraph.text.replace(placeholder, value_str)
                                shape.text_frame.word_wrap = True
                                self._auto_fit_text(paragraph, shape.width, shape.height)
                                count += 1
        logger.info(f"      共替换 {count} 个文本占位符")
        return count
    
    def _auto_fit_text(self, paragraph, box_width, box_height, min_font_size=10, max_font_size=14):
        """自动调整字体大小以适应文本框"""
        from pptx.util import Pt
        text_length = len(paragraph.text)
        box_width_inches = box_width.inches
        chars_per_line = int(box_width_inches * max_font_size * 1.2)
        needed_lines = text_length / chars_per_line + 1
        box_height_inches = box_height.inches
        max_lines = int(box_height_inches * 15)
        if needed_lines > max_lines:
            font_size = max(min_font_size, int(max_font_size * max_lines / needed_lines))
            paragraph.font.size = Pt(font_size)
    
    def find_placeholder_positions(self, prs: Presentation, placeholder_pattern: str) -> List[Dict]:
        """查找占位符位置（兼容单方括号和双方括号）"""
        import re
        positions = []
        
        # 提取图表名称
        chart_name = None
        if ':' in placeholder_pattern:
            # 处理 [CHART:xxx] 或 [[CHART:xxx]]
            clean_pattern = placeholder_pattern.strip('[]')
            chart_name = clean_pattern.split(':')[1] if ':' in clean_pattern else None
        
        # 构建正则表达式（兼容单方括号和双方括号）
        # 例如：[CHART:xxx] 或 [[CHART:xxx]]
        if chart_name:
            regex_pattern = rf'\[\[?CHART:{re.escape(chart_name)}\]\]?'
        else:
            regex_pattern = re.escape(placeholder_pattern)
        
        regex = re.compile(regex_pattern)
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # 检查 shape 名称
                if chart_name and hasattr(shape, 'name') and shape.name == f"CHART:{chart_name}":
                    positions.append({
                        'slide_index': slide_idx, 'shape_index': shape_idx,
                        'placeholder': placeholder_pattern, 'left': shape.left,
                        'top': shape.top, 'width': shape.width, 'height': shape.height,
                        'shape': shape, 'slide': slide, 'by_name': True
                    })
                    continue
                
                # 检查文本框内容
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text
                        if regex.search(text):
                            positions.append({
                                'slide_index': slide_idx, 'shape_index': shape_idx,
                                'placeholder': placeholder_pattern, 'left': shape.left,
                                'top': shape.top, 'width': shape.width, 'height': shape.height,
                                'shape': shape, 'slide': slide, 'by_name': False
                            })
        return positions
    
    def replace_with_chart(self, prs: Presentation, placeholder: str, chart_path: str) -> bool:
        """替换图表占位符"""
        positions = self.find_placeholder_positions(prs, placeholder)
        if not positions:
            # 静默失败（可能是因为另一种括号格式已经成功替换）
            return False
        
        replaced_count = 0
        for pos in positions:
            slide = pos['slide']
            shape = pos['shape']
            try:
                sp = shape.element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(chart_path, pos['left'], pos['top'], width=pos['width'], height=pos['height'])
                replaced_count += 1
                logger.info(f"      已替换图表 {placeholder} (slide {pos['slide_index']+1})")
            except Exception as e:
                logger.error(f"替换图表失败：{e}")
        return replaced_count > 0
    
    def save(self, prs: Presentation, output_path: str):
        """保存 PPT"""
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        logger.info(f"      PPT 已保存：{output_path}")
    
    def get_placeholder_config(self) -> Dict:
        """获取占位符配置"""
        return self.placeholders
    
    def export_placeholders(self, output_path: str = None):
        """导出占位符配置到文件"""
        if output_path is None:
            output_path = self.placeholders_file
        self._save_placeholders(self.placeholders)
        logger.info(f"占位符配置已导出：{output_path}")
    
    def regenerate_placeholders_from_template(self, template_name: str):
        """从模板重新生成占位符配置"""
        template_path = os.path.join(self.templates_dir, template_name)
        if not os.path.exists(template_path):
            logger.error(f"模板不存在：{template_path}")
            return False
        
        self.placeholders = self.scan_placeholders_from_template(template_path)
        self._save_placeholders(self.placeholders)
        logger.info(f"占位符配置已重新生成：{self.placeholders_file}")
        return True
