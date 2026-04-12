# -*- coding: utf-8 -*-
"""
PPT 原生图表引擎 - 在 PPT 中创建可编辑的图表
使用 python-pptx 的 Chart 功能
"""
from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd
from typing import Dict, List, Optional, Tuple
import logging

logger = logging.getLogger(__name__)


class PPTChartEngine:
    """PPT 原生图表引擎 - 创建可编辑的图表"""
    
    def __init__(self):
        # 图表类型映射
        self.chart_type_map = {
            'bar_horizontal': XL_CHART_TYPE.BAR_CLUSTERED,
            'bar_vertical': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'pie': XL_CHART_TYPE.PIE,
            'line': XL_CHART_TYPE.LINE,
            'column_clustered': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'area': XL_CHART_TYPE.AREA,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
        }
    
    def create_chart_in_placeholder(self, 
                                     placeholder, 
                                     df: pd.DataFrame, 
                                     chart_type: str,
                                     title: str = '',
                                     **kwargs) -> bool:
        """
        在 PPT 占位符中创建原生图表
        
        Args:
            placeholder: PPT 占位符对象
            df: 数据 DataFrame
            chart_type: 图表类型
            title: 图表标题
            **kwargs: 其他参数（x_field, y_field, category_field, series 等）
        
        Returns:
            bool: 是否成功创建
        """
        try:
            # 获取占位符的位置和大小
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height
            
            # 删除占位符
            sp = placeholder.element
            sp.getparent().remove(sp)
            
            # 获取幻灯片对象
            slide = placeholder.parent
            
            # 根据图表类型创建不同的图表
            if chart_type in ['bar_horizontal', 'bar_vertical', 'column_clustered']:
                return self._create_categorical_chart(slide, df, chart_type, title, left, top, width, height, **kwargs)
            elif chart_type == 'pie':
                return self._create_pie_chart(slide, df, title, left, top, width, height, **kwargs)
            elif chart_type == 'line':
                return self._create_line_chart(slide, df, title, left, top, width, height, **kwargs)
            elif chart_type == 'area':
                return self._create_area_chart(slide, df, title, left, top, width, height, **kwargs)
            elif chart_type == 'scatter':
                return self._create_scatter_chart(slide, df, title, left, top, width, height, **kwargs)
            else:
                logger.warning(f"不支持的原生图表类型：{chart_type}")
                return False
                
        except Exception as e:
            logger.error(f"创建原生图表失败：{e}")
            return False
    
    def _create_categorical_chart(self, slide, df: pd.DataFrame, chart_type: str, title: str,
                                   left, top, width, height, **kwargs):
        """创建分类图表（柱状图/条形图）"""
        category_field = kwargs.get('category_field')
        series = kwargs.get('series', [])
        x_field = kwargs.get('x_field')
        y_field = kwargs.get('y_field')
        
        # 准备数据
        chart_data = CategoryChartData()
        
        # 设置分类
        if category_field and category_field in df.columns:
            chart_data.categories = df[category_field].tolist()
        elif y_field and y_field in df.columns:
            chart_data.categories = df[y_field].tolist()
        
        # 添加系列
        if series and isinstance(series, list):
            # 多系列
            for s in series:
                if s in df.columns:
                    chart_data.add_series(s, df[s].tolist())
        elif x_field and y_field:
            # 单系列
            if x_field in df.columns and y_field in df.columns:
                chart_data.add_series(x_field, df[y_field].tolist())
        
        # 创建图表
        chart_type_enum = self.chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart = slide.shapes.add_chart(chart_type_enum, left, top, width, height, chart_data).chart
        
        # 设置标题
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        # 设置图例
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        logger.info(f"已创建原生分类图表：{title}")
        return True
    
    def _create_pie_chart(self, slide, df: pd.DataFrame, title: str, left, top, width, height, **kwargs):
        """创建饼图"""
        category_field = kwargs.get('category_field')
        value_field = kwargs.get('value_field')
        
        if not category_field or not value_field:
            logger.error("饼图需要 category_field 和 value_field")
            return False
        
        if category_field not in df.columns or value_field not in df.columns:
            logger.error(f"字段不存在：category_field={category_field}, value_field={value_field}")
            return False
        
        # 准备数据
        chart_data = CategoryChartData()
        chart_data.categories = df[category_field].tolist()
        chart_data.add_series(value_field, df[value_field].tolist())
        
        # 创建图表
        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, chart_data).chart
        
        # 设置标题
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        # 设置数据标签（显示百分比）
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = True
        data_labels.show_category_name = True
        
        logger.info(f"已创建原生饼图：{title}")
        return True
    
    def _create_line_chart(self, slide, df: pd.DataFrame, title: str, left, top, width, height, **kwargs):
        """创建折线图"""
        x_field = kwargs.get('x_field')
        y_field = kwargs.get('y_field')
        
        if not x_field or not y_field:
            logger.error("折线图需要 x_field 和 y_field")
            return False
        
        if x_field not in df.columns or y_field not in df.columns:
            logger.error(f"字段不存在：x_field={x_field}, y_field={y_field}")
            return False
        
        # 准备数据
        chart_data = CategoryChartData()
        chart_data.categories = df[x_field].tolist()
        chart_data.add_series(y_field, df[y_field].tolist())
        
        # 创建图表
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, left, top, width, height, chart_data).chart
        
        # 设置标题
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        # 设置图例
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        logger.info(f"已创建原生折线图：{title}")
        return True
    
    def _create_area_chart(self, slide, df: pd.DataFrame, title: str, left, top, width, height, **kwargs):
        """创建面积图"""
        x_field = kwargs.get('x_field')
        y_fields = kwargs.get('y_fields', [])
        
        if not x_field:
            logger.error("面积图需要 x_field")
            return False
        
        if x_field not in df.columns:
            logger.error(f"字段不存在：x_field={x_field}")
            return False
        
        # 准备数据
        chart_data = CategoryChartData()
        chart_data.categories = df[x_field].tolist()
        
        # 添加系列
        if y_fields and isinstance(y_fields, list):
            for y_field in y_fields:
                if y_field in df.columns:
                    chart_data.add_series(y_field, df[y_field].tolist())
        elif len(df.columns) > 1:
            # 自动添加所有数值列
            for col in df.columns:
                if col != x_field and pd.api.types.is_numeric_dtype(df[col]):
                    chart_data.add_series(col, df[col].tolist())
        
        # 创建图表
        chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA, left, top, width, height, chart_data).chart
        
        # 设置标题
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        logger.info(f"已创建原生面积图：{title}")
        return True
    
    def _create_scatter_chart(self, slide, df: pd.DataFrame, title: str, left, top, width, height, **kwargs):
        """创建散点图"""
        x_field = kwargs.get('x_field')
        y_field = kwargs.get('y_field')
        
        if not x_field or not y_field:
            logger.error("散点图需要 x_field 和 y_field")
            return False
        
        if x_field not in df.columns or y_field not in df.columns:
            logger.error(f"字段不存在：x_field={x_field}, y_field={y_field}")
            return False
        
        # 准备数据
        chart_data = XyChartData()
        
        # 添加系列
        series = chart_data.add_series("数据点")
        for _, row in df.iterrows():
            x_val = row[x_field]
            y_val = row[y_field]
            # 尝试转换为数值
            try:
                x_val = float(x_val)
                y_val = float(y_val)
                series.add_data_point(x_val, y_val)
            except:
                continue
        
        # 创建图表
        chart = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, left, top, width, height, chart_data).chart
        
        # 设置标题
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        logger.info(f"已创建原生散点图：{title}")
        return True
